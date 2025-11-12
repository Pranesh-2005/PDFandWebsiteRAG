# app.py
import os
import asyncio
import json
import hashlib
import shutil
from io import BytesIO, StringIO
from typing import List, Tuple

import gradio as gr
import numpy as np
import faiss
import requests
import pandas as pd
from sentence_transformers import SentenceTransformer
import fitz  # PyMuPDF
import docx
from pptx import Presentation
from crawl4ai import AsyncWebCrawler

# ---------------- Config ----------------
OPENROUTER_API_KEY = os.getenv("OPENROUTER_API_KEY")
OPENROUTER_MODEL = "nvidia/nemotron-nano-12b-v2-vl:free"
EMBEDDING_MODEL_NAME = "all-MiniLM-L6-v2"
CACHE_DIR = "./cache"
SYSTEM_PROMPT = "You are a helpful assistant."
os.makedirs(CACHE_DIR, exist_ok=True)

embedder = SentenceTransformer(EMBEDDING_MODEL_NAME)

DOCS: List[str] = []
FILENAMES: List[str] = []
EMBEDDINGS: np.ndarray = None
FAISS_INDEX = None
CURRENT_CACHE_KEY: str = ""


# ---------------- Periodic cache cleanup ----------------
async def clear_cache_every_5min():
    while True:
        await asyncio.sleep(300)  # 5 minutes
        try:
            if os.path.exists(CACHE_DIR):
                shutil.rmtree(CACHE_DIR)
            os.makedirs(CACHE_DIR, exist_ok=True)
            print("🧹 Cache cleared successfully.")
        except Exception as e:
            print(f"[Cache cleanup error] {e}")

# Launch the cleaner in background
asyncio.get_event_loop().create_task(clear_cache_every_5min())

def extract_text_from_pdf(file_bytes: bytes) -> str:
    try:
        doc = fitz.open(stream=file_bytes, filetype="pdf")
        return "\n".join(page.get_text() for page in doc)
    except Exception as e:
        return f"[PDF extraction error] {e}"

def extract_text_from_docx(file_bytes: bytes) -> str:
    try:
        f = BytesIO(file_bytes)
        doc = docx.Document(f)
        return "\n".join(p.text for p in doc.paragraphs)
    except Exception as e:
        return f"[DOCX extraction error] {e}"

def extract_text_from_txt(file_bytes: bytes) -> str:
    try:
        return file_bytes.decode("utf-8", errors="ignore")
    except Exception as e:
        return f"[TXT extraction error] {e}"

def extract_text_from_excel(file_bytes: bytes) -> str:
    try:
        f = BytesIO(file_bytes)
        df = pd.read_excel(f, dtype=str)
        return "\n".join("\n".join(df[col].fillna("").astype(str).tolist()) for col in df.columns)
    except Exception as e:
        return f"[EXCEL extraction error] {e}"

def extract_text_from_pptx(file_bytes: bytes) -> str:
    try:
        f = BytesIO(file_bytes)
        prs = Presentation(f)
        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    texts.append(shape.text)
        return "\n".join(texts)
    except Exception as e:
        return f"[PPTX extraction error] {e}"

def extract_text_from_csv(file_bytes: bytes) -> str:
    try:
        f = StringIO(file_bytes.decode("utf-8", errors="ignore"))
        df = pd.read_csv(f, dtype=str)
        return df.to_string(index=False)
    except Exception as e:
        return f"[CSV extraction error] {e}"

def extract_text_from_file_tuple(file_tuple) -> Tuple[str, bytes]:
    try:
        if hasattr(file_tuple, "name") and hasattr(file_tuple, "read"):
            return os.path.basename(file_tuple.name), file_tuple.read()
    except Exception:
        pass
    if isinstance(file_tuple, tuple) and len(file_tuple) == 2 and isinstance(file_tuple[1], (bytes, bytearray)):
        return file_tuple[0], bytes(file_tuple[1])
    if isinstance(file_tuple, str) and os.path.exists(file_tuple):
        with open(file_tuple, "rb") as fh:
            return os.path.basename(file_tuple), fh.read()
    raise ValueError("Unsupported file object passed by Gradio.")

def extract_text_by_ext(filename: str, file_bytes: bytes) -> str:
    name = filename.lower()
    if name.endswith(".pdf"): return extract_text_from_pdf(file_bytes)
    if name.endswith(".docx"): return extract_text_from_docx(file_bytes)
    if name.endswith(".txt"): return extract_text_from_txt(file_bytes)
    if name.endswith((".xlsx", ".xls")): return extract_text_from_excel(file_bytes)
    if name.endswith(".pptx"): return extract_text_from_pptx(file_bytes)
    if name.endswith(".csv"): return extract_text_from_csv(file_bytes)
    return extract_text_from_txt(file_bytes)



# ---------------- Cache + FAISS helpers ----------------
def make_cache_key_for_files(files: List[Tuple[str, bytes]]) -> str:
    h = hashlib.sha256()
    for name, b in sorted(files, key=lambda x: x[0]):
        h.update(name.encode())
        h.update(str(len(b)).encode())
        h.update(hashlib.sha256(b).digest())
    return h.hexdigest()

def cache_save_embeddings(cache_key: str, embeddings: np.ndarray, filenames: List[str]):
    np.savez_compressed(os.path.join(CACHE_DIR, f"{cache_key}.npz"), embeddings=embeddings, filenames=np.array(filenames))

def cache_load_embeddings(cache_key: str):
    path = os.path.join(CACHE_DIR, f"{cache_key}.npz")
    if not os.path.exists(path): return None
    try:
        arr = np.load(path, allow_pickle=True)
        return arr["embeddings"], arr["filenames"].tolist()
    except Exception:
        return None

def build_faiss_index(embeddings: np.ndarray):
    global FAISS_INDEX
    if embeddings is None or len(embeddings) == 0:
        FAISS_INDEX = None
        return None
    emb = embeddings.astype("float32")
    index = faiss.IndexFlatL2(emb.shape[1])
    index.add(emb)
    FAISS_INDEX = index
    return index

def search_top_k(query: str, k: int = 3):
    if FAISS_INDEX is None:
        return []
    q_emb = embedder.encode([query], convert_to_numpy=True).astype("float32")
    D, I = FAISS_INDEX.search(q_emb, k)
    return [{"index": int(i), "distance": float(d), "text": DOCS[i], "source": FILENAMES[i]} for d, i in zip(D[0], I[0]) if i >= 0]


# ---------------- OpenRouter Client ----------------
def openrouter_chat_system_user(user_prompt: str):
    """
    Sends user prompt to OpenRouter and expects a plain text response.
    """
    if not OPENROUTER_API_KEY:
        return "[OpenRouter error] Missing OPENROUTER_API_KEY."

    url = "https://openrouter.ai/api/v1/chat/completions"
    headers = {
        "Authorization": f"Bearer {OPENROUTER_API_KEY}",
        "Content-Type": "application/json",
    }

    # Tell the model explicitly to reply as plain text only
    payload = {
        "model": OPENROUTER_MODEL,
        "messages": [
            {"role": "system", "content": SYSTEM_PROMPT + " Always respond in plain text. Avoid JSON or markdown formatting."},
            {"role": "user", "content": user_prompt},
        ],
    }

    try:
        r = requests.post(url, headers=headers, json=payload, timeout=60)
        r.raise_for_status()
        obj = r.json()

        # Safely extract plain text
        if "choices" in obj and obj["choices"]:
            choice = obj["choices"][0]
            if "message" in choice and "content" in choice["message"]:
                text = choice["message"]["content"]
                # Ensure no markdown or code blocks
                text = text.strip().replace("```", "").replace("json", "")
                return text
            elif "text" in choice:
                return choice["text"].strip()
        return "[OpenRouter] Unexpected response format."

    except Exception as e:
        return f"[OpenRouter request error] {e}"


# ---------------- Crawl4AI Logic ----------------
async def _crawl_async_get_markdown(url: str):
    async with AsyncWebCrawler() as crawler:
        result = await crawler.arun(url=url)
        if hasattr(result, "success") and result.success is False:
            return f"[Crawl4AI error] {getattr(result, 'error_message', '[Unknown error]')}"
        md_obj = getattr(result, "markdown", None)
        if md_obj:
            return getattr(md_obj, "fit_markdown", None) or getattr(md_obj, "raw_markdown", None) or str(md_obj)
        return getattr(result, "text", None) or getattr(result, "html", None) or "[Crawl4AI returned no usable fields]"

def crawl_url_sync(url: str) -> str:
    try:
        return asyncio.run(_crawl_async_get_markdown(url))
    except Exception as e:
        return f"[Crawl4AI runtime error] {e}"


# ---------------- Gradio Handlers ----------------
def upload_and_index(files):
    global DOCS, FILENAMES, EMBEDDINGS, CURRENT_CACHE_KEY
    if not files:
        return "No files uploaded.", ""
    prepared = [(name := extract_text_from_file_tuple(f)[0], extract_text_from_file_tuple(f)[1]) for f in files]
    previews = [{"name": n, "size": len(b)} for n, b in prepared]
    cache_key = make_cache_key_for_files(prepared)
    CURRENT_CACHE_KEY = cache_key
    cached = cache_load_embeddings(cache_key)
    if cached:
        emb, filenames = cached
        EMBEDDINGS = np.array(emb)
        FILENAMES = filenames
        DOCS = [extract_text_by_ext(n, b) for n, b in prepared]
        build_faiss_index(EMBEDDINGS)
        return f"Loaded embeddings from cache ({len(FILENAMES)} docs).", json.dumps(previews)
    DOCS, FILENAMES = zip(*[(extract_text_by_ext(n, b), n) for n, b in prepared])
    EMBEDDINGS = embedder.encode(DOCS, convert_to_numpy=True, show_progress_bar=False).astype("float32")
    cache_save_embeddings(cache_key, EMBEDDINGS, FILENAMES)
    build_faiss_index(EMBEDDINGS)
    return f"Uploaded and indexed {len(DOCS)} documents.", json.dumps(previews)

def crawl_and_index(url: str):
    global DOCS, FILENAMES, EMBEDDINGS, CURRENT_CACHE_KEY
    if not url:
        return "No URL provided.", ""
    crawled = crawl_url_sync(url)
    if crawled.startswith("[Crawl4AI"):
        return crawled, ""
    key_hash = hashlib.sha256((url + crawled).encode()).hexdigest()
    CURRENT_CACHE_KEY = key_hash
    cached = cache_load_embeddings(key_hash)
    if cached:
        emb, filenames = cached
        EMBEDDINGS = np.array(emb)
        FILENAMES = filenames
        DOCS = [crawled]
        build_faiss_index(EMBEDDINGS)
        return f"Crawled and loaded embeddings from cache for {url}", crawled[:20000]
    DOCS, FILENAMES = [crawled], [url]
    EMBEDDINGS = embedder.encode(DOCS, convert_to_numpy=True, show_progress_bar=False).astype("float32")
    cache_save_embeddings(key_hash, EMBEDDINGS, FILENAMES)
    build_faiss_index(EMBEDDINGS)
    return f"Crawled and indexed {url}", crawled[:20000]

def ask_question(question: str):
    if not question:
        return "Please enter a question."
    if not DOCS or FAISS_INDEX is None:
        return "No indexed data found."
    results = search_top_k(question, k=3)
    if not results:
        return "No relevant documents found."
    context = "\n".join(f"Source: {r['source']}\n\n{r['text'][:18000]}\n---\n" for r in results)
    user_prompt = f"Use the following context to answer the question.\n\nContext:\n{context}\nQuestion: {question}\nAnswer:"
    return openrouter_chat_system_user(user_prompt)


# ---------------- Gradio UI ----------------
with gr.Blocks(title="AI Ally — Crawl4AI + OpenRouter + FAISS") as demo:
    gr.Markdown("# 🤖 AI Ally — Document & Website QA\nCrawl4AI for websites, file uploads for docs. FAISS retrieval + sentence-transformers + OpenRouter LLM.")

    with gr.Tab("Documents"):
        file_input = gr.File(label="Upload files", file_count="multiple",
                             file_types=[".pdf", ".docx", ".txt", ".xlsx", ".pptx", ".csv"])
        upload_btn = gr.Button("Upload & Index")
        upload_status = gr.Textbox(label="Status", interactive=False)
        preview_box = gr.Textbox(label="Uploads (preview JSON)", interactive=False)
        upload_btn.click(upload_and_index, inputs=[file_input], outputs=[upload_status, preview_box])

        gr.Markdown("### Ask about your documents")
        q = gr.Textbox(label="Question", lines=3)
        ask_btn = gr.Button("Ask")
        answer_out = gr.Textbox(label="Answer", interactive=False, lines=15)
        ask_btn.click(ask_question, inputs=[q], outputs=[answer_out])

    with gr.Tab("Website Crawl"):
        url = gr.Textbox(label="URL to crawl")
        crawl_btn = gr.Button("Crawl & Index")
        crawl_status = gr.Textbox(label="Status", interactive=False)
        crawl_preview = gr.Textbox(label="Crawl preview", interactive=False)
        crawl_btn.click(crawl_and_index, inputs=[url], outputs=[crawl_status, crawl_preview])

        q2 = gr.Textbox(label="Question", lines=3)
        ask_btn2 = gr.Button("Ask site")
        answer_out2 = gr.Textbox(label="Answer", interactive=False, lines=15)
        ask_btn2.click(ask_question, inputs=[q2], outputs=[answer_out2])

    with gr.Tab("Settings / Info"):
        gr.Markdown(f"- Model: `{OPENROUTER_MODEL}`")
        gr.Markdown(f"- Embedding model: `{EMBEDDING_MODEL_NAME}`")
        gr.Markdown(f"- Cache clears automatically every 5 minutes.")
        gr.Markdown(f"- System prompt is fixed internally: `{SYSTEM_PROMPT}`")

if __name__ == "__main__":
    demo.launch(server_name="0.0.0.0", server_port=7860, debug=True)